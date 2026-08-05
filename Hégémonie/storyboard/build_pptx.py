#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Construit les PowerPoint de table a partir des images de out/.

Chaque diapo = une image plein cadre, sans titre ni habillage.
Les notes du presentateur contiennent : quand montrer, quoi lire, quoi
distribuer, et les 3 pistes musicales de la scene.

Deux fichiers sont produits :
  Le_Service_du_soir.pptx               leger, musique en playlist a part
  Le_Service_du_soir_AVEC_MUSIQUE.pptx  mp3 embarques, lecture au clic

Depose des mp3 dans ce dossier et relance : ils sont apparies automatiquement
au mot-cle de chaque piste, et une icone apparait par piste trouvee.

Usage :
    python build_pptx.py
"""

import unicodedata
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
OUT = HERE / "out"
SLIDES = OUT / "slides"

SLIDE_W, SLIDE_H = 13.333, 7.5          # 16:9
TARGET_PX = 1920                         # largeur des images embarquees

OR = RGBColor(0xE8, 0xC9, 0x8A)
GRIS = RGBColor(0x9A, 0x8A, 0x76)
NOIR = RGBColor(0x0D, 0x0B, 0x09)

# ---------------------------------------------------------------- musique
# 3 pistes par scene. Detail et justifications : CONDUITE_MUSICALE.md
# Chaque piste : (titre, source, duree approx, mot-cle d'appariement du mp3)

MUSIQUE = {
    1: {"cue": "Avant ta première phrase. COUPE NET quand la femelle décolle.",
        "pistes": [
            ("Ambre", "Nils Frahm — Felt", "~4:00", "ambre"),
            ("Subwoofer Lullaby", "C418 — Minecraft", "~3:30", "subwoofer"),
            ("Only the Winds", "Ólafur Arnalds", "~4:00", "only the winds"),
        ]},
    2: {"cue": "Sur « La femelle ne touche plus le sol. » Pas avant.",
        "pistes": [
            ("Flight From the City", "Jóhann Jóhannsson — Orphée", "~6:20", "flight from the city"),
            ("Nascence", "Austin Wintory — Journey", "~2:30", "nascence"),
            ("Adventure", "Disasterpeace — Fez", "~3:00", "adventure"),
        ]},
    3: {"cue": "Sur « Une ombre passe sous la coque. » LAISSE-LA FINIR dans le noir.",
        "pistes": [
            ("An Ending (Ascent)", "Brian Eno — Apollo", "~4:25", "an ending"),
            ("Cetus", "Austin Wintory — Abzû", "~4:00", "cetus"),
            ("Travelers", "Andrew Prahlow — Outer Wilds", "~3:00", "travelers"),
        ]},
    4: {"cue": "Juste avant « Vous vous attendiez à du noir. » Laisse respirer plus longtemps.",
        "pistes": [
            ("Says", "Nils Frahm — Spaces", "~8:00", "says"),
            ("Dirtmouth", "Christopher Larkin — Hollow Knight", "~2:30", "dirtmouth"),
            ("Delphinus Delphis", "Austin Wintory — Abzû", "~3:30", "delphinus"),
        ]},
    5: {"cue": "Dès le quartier commerçant. BAISSE DE MOITIÉ quand Bumbur fait son signe.",
        "pistes": [
            ("Days to Come", "Bonobo", "~4:00", "bonobo"),
            ("Spiritfarer (Main Theme)", "Max LL — Spiritfarer", "~3:00", "spiritfarer"),
            ("Whirling-In-Rags, 8 AM", "British Sea Power — Disco Elysium", "~4:00", "whirling"),
        ]},
    6: {"cue": "En entrant dans la salle des roues. Coupe quand la négociation commence.",
        "pistes": [
            ("The Cause of Labour…", "Jóhann Jóhannsson — The Miners' Hymns", "~6:00", "cause of labour"),
            ("The City Must Survive", "Piotr Musiał — Frostpunk", "~4:00", "city must survive"),
            ("They Being Dead Yet Speaketh", "Jóhann Jóhannsson", "~5:00", "they being dead"),
        ]},
    7: {"cue": "Au moment où tu avances le marqueur d'une zone. Sans commentaire.",
        "pistes": [
            ("On the Nature of Daylight", "Max Richter", "~6:10", "nature of daylight"),
            ("Prologue ~ To the Ancient Land", "Kow Otani — Shadow of the Colossus", "~4:00", "ancient land"),
            ("The Sacrifice", "Gareth Coker — Ori", "~3:00", "sacrifice"),
        ]},
    8: {"cue": "Quand Pagure annonce l'ordre du jour. COUPE BRUTALEMENT si un joueur nomme la cathédrale.",
        "pistes": [
            ("The Sun's Gone Dim…", "Jóhann Jóhannsson — IBM 1401", "~5:00", "suns gone dim"),
            ("Papers, Please (Theme)", "Lucas Pope", "~2:30", "papers please"),
            ("Resting Grounds", "Christopher Larkin — Hollow Knight", "~3:00", "resting grounds"),
        ]},
    9: {"cue": "ATTENTION : les 4 minutes de silence se jouent SANS musique. Lance quand la minuterie sonne.",
        "pistes": [
            ("Apotheosis", "Austin Wintory — Journey", "~3:40", "apotheosis"),
            ("Light of Nibel", "Gareth Coker — Ori", "~4:00", "light of nibel"),
            ("BB's Theme", "Ludvig Forssell — Death Stranding", "~4:00", "bb"),
        ]},
    10: {"cue": "Sur la dernière image. Laisse-la FINIR avant de poser la question à Bumbur.",
         "pistes": [
             ("Near Light", "Ólafur Arnalds — Living Room Songs", "~3:10", "near light"),
             ("Sweden", "C418 — Minecraft", "~3:30", "sweden"),
             ("Hateno Village", "Manaka Kataoka — Breath of the Wild", "~3:00", "hateno"),
         ]},
}

# ---------------------------------------------------------------- contenu

TITRE = {
    "titre": "LE SERVICE DU SOIR",
    "sous": "Hégémonie — one-shot",
    "notes": (
        "ÉCRAN NOIR AVANT DE COMMENCER : touche B.\n"
        "Le noir est ton état par défaut. Une image ne reste jamais affichée "
        "pendant qu'on joue — tu la montres, tu laisses 10 secondes, tu reprends "
        "le noir. Sinon les joueurs regardent l'écran au lieu de se parler.\n\n"
        "♪ MUSIQUE : {MUSIQUE_INTRO}\n"
        "Trois pistes par scène, à enchaîner dans l'ordre : ~10 à 15 min de musique "
        "par scène, environ 2 h pour la séance.\n"
        "Deux moments sans musique du tout : la négociation avec les marmottes une "
        "fois entamée, et les 4 minutes de silence de l'acte III.\n\n"
        "Rappel de ton : personne ne dit jamais « on est dans une baleine ». "
        "Les habitants disent « on habite ici »."
    ),
}

NOTES = {
    1: {
        "titre": "Les dernières réserves",
        "quand": "ACTE I — scène 1, dès l'ouverture.",
        "lire": (
            "Ça fait quatre mois. Quatre mois que les cristaux sont morts, quatre mois "
            "qu'il n'y a pas un souffle, quatre mois que le navire est exactement au même "
            "endroit du vide. Vous avez déposé Brime, Prune et Lyra à la dernière île "
            "habitée — c'était plus raisonnable.\n"
            "Il ne reste plus rien dans la cuisine. Rien. Vous avez fini le sel.\n"
            "Il reste deux moutons dans la cale. Un mâle et une femelle. Vous les gardiez "
            "pour les faire s'accoupler, parce que c'était la seule idée d'avenir que vous "
            "ayez eue depuis quatre mois.\n"
            "Et ce matin, quelqu'un a dit à voix haute qu'il faudrait peut-être en manger un."
        ),
        "apres": "Laisse-les discuter cinq minutes. Ne mentionne jamais Bumbur.",
    },
    2: {
        "titre": "Le mouton s'envole",
        "quand": "ACTE I — scène 2.",
        "lire": (
            "La femelle ne touche plus le sol. Dix centimètres. Elle a l'air contente. "
            "Elle vous regarde comme si vous alliez faire quelque chose.\n"
            "Le temps que quelqu'un monte chercher une corde, elle fait la taille d'un veau. "
            "Le temps de redescendre, elle a la taille d'une vache et elle vient de passer "
            "à travers une cloison — pas violemment. Juste en grandissant au travers.\n"
            "Et le bêlement descend. À chaque fois qu'elle ouvre la bouche, c'est plus grave "
            "que la fois d'avant.\n"
            "Le mâle, lui, ne bouge pas. Il la regarde."
        ),
        "apres": "Enchaîne vite sur la scène 3. Ne laisse pas retomber.",
    },
    3: {
        "titre": "La baleine astrale",
        "quand": "ACTE I — scène 3, l'événement déclencheur.",
        "lire": (
            "Une ombre passe sous la coque. Puis elle continue de passer. Pendant onze "
            "secondes.\n"
            "Il n'y a aucune agressivité là-dedans. Aucune. C'est le geste de quelqu'un qui "
            "ramasse un enfant tombé dans la rue — et vous êtes ce qu'il y avait autour.\n"
            "Noir. Un bruit de porte. Le silence.\n"
            "Puis, trente secondes plus tard, quelque part au-dessus de vous : une cloche. "
            "Deux coups. Et une voix, très loin, qui gueule « c'est l'heure ! »"
        ),
        "apres": "AUCUN JET. Aucune échappatoire. Puis écran noir, et pause de 3 secondes.",
    },
    4: {
        "titre": "La ville dans la baleine",
        "quand": "ACTE II — 2A, l'arrivée. LA grande image de la soirée.",
        "lire": (
            "Vous vous attendiez à du noir, de l'acide, une odeur.\n"
            "Il y a du linge qui sèche.\n"
            "Une quarantaine de bateaux amarrés les uns aux autres, des passerelles entre "
            "les mâts, des enseignes peintes, des gens qui portent des caisses. Quelqu'un "
            "râle sur le prix de quelque chose. Il fait tiède.\n"
            "Un bernard-l'ermite de la taille d'un homme, calé dans une coquille scellée au "
            "pont, vous regarde arriver avec exactement l'expression de quelqu'un qui voit "
            "une voiture se garer sur sa place."
        ),
        "apres": "Laisse celle-ci affichée plus longtemps. Pose le marqueur sur la zone 6.",
    },
    5: {
        "titre": "À la Bonne Étoile",
        "quand": "ACTE II — 2B, les retrouvailles avec Bumbur.",
        "lire": (
            "L'enseigne dit À LA BONNE ÉTOILE, avec une étoile peinte à côté.\n"
            "La salle est pleine. Il fait chaud, ça sent le beurre. Un poulpe en gilet place "
            "les clients — sept bras, pas huit. Un autre verse le vin. Un troisième "
            "débarrasse. Il n'y a que des poulpes.\n"
            "Et au fond, derrière le passe, penché sur une casserole, un halfelin "
            "rondouillard vous voit, s'arrête net — et vous fait un grand signe de la main. "
            "Un vrai. Il a l'air sincèrement, complètement heureux.\n"
            "Puis il crie « deux mâts, une soupe ! » et il retourne à ses casseroles."
        ),
        "apres": (
            "Distribue le menu imprimé MAINTENANT.\n"
            "Il les nourrit avant toute conversation, sans réaliser que c'est un événement.\n"
            "Personne dans ce quartier n'a jamais vu d'étoile — ils ne comprennent pas le nom."
        ),
    },
    6: {
        "titre": "La centrale des marmottes",
        "quand": "ACTE II — 2C.",
        "lire": (
            "Une salle grande comme une cathédrale — une vraie, pas celle en os. Des "
            "centaines de roues sur trois niveaux, des courroies, des tableaux de report "
            "avec toutes les aiguilles à zéro.\n"
            "Et quatre cents marmottes assises par terre, en silence, qui vous regardent "
            "entrer."
        ),
        "apres": (
            "Motmot : « Le quartier s'est fait avant les cristaux. C'est le vieux système. »\n"
            "Sors le PRÉAVIS manuscrit. Elles ne savent ni lire ni écrire.\n"
            "Le contrat imprimé est l'accessoire le plus fort de la soirée — garde-le pour "
            "quand ils poussent."
        ),
    },
    7: {
        "titre": "Le veau qui écrase le quartier",
        "quand": "ACTE II — au 1er ou 2e palier de l'horloge, quand ça devient sérieux.",
        "lire": (
            "Il est trop grand pour la rue. Il ne fait rien de méchant : il respire, et à "
            "chaque fois qu'il respire une passerelle cède.\n"
            "Il a de la laine sur le dos et une oreille de mouton. Ses yeux ont commencé à "
            "migrer sur les côtés de sa tête et se sont arrêtés à mi-chemin.\n"
            "Il a peur. Et il grandit toujours."
        ),
        "apres": (
            "Avance le marqueur d'une zone à vue de tous, sans commenter.\n"
            "L'attaquer le fait grandir plus vite : +1 palier. Le calmer : −1."
        ),
    },
    8: {
        "titre": "L'Assemblée générale",
        "quand": "ACTE II — quand la zone 5 tombe. LA scène centrale.",
        "lire": (
            "Tout le monde est là. Les propriétaires d'un côté, scellés au pont dans leurs "
            "coquilles, un bulletin dans les pinces. Les locataires de l'autre, onze mille "
            "sardines dans leurs sphères. Les marmottes debout au fond, qui n'ont pas le "
            "droit de vote.\n"
            "Et au-dessus de tout le monde, accrochés aux côtes : les harpons.\n"
            "MOTION 1 — abattage de la bête, selon la procédure de 1966.\n"
            "MOTION 2 — mise en cause des occupants du lot 14."
        ),
        "apres": (
            "DISTRIBUE LE BULLETIN DE VOTE aux joueurs — ils n'ont pas le droit de vote, ils "
            "sont l'objet de la motion 2. Les regarder comprendre, c'est la moitié de la "
            "scène.\n"
            "Pagure : « Il faut un responsable. Sinon personne ne lèvera la main. »\n"
            "Si un joueur nomme la cathédrale — « vous l'avez déjà fait, regardez où vous "
            "priez » — LA SALLE SE TAIT."
        ),
    },
    9: {
        "titre": "Le premier chant",
        "quand": "ACTE III — le climax. Après les 4 minutes de silence.",
        "lire": (
            "MONTRE L'IMAGE, PUIS NE DIS PLUS RIEN.\n"
            "Minuterie sur 4:00. Plus personne ne parle — ni les joueurs, ni toi.\n"
            "Tu décris en écrivant sur un papier que tu fais passer."
        ),
        "apres": (
            "Pendant ces 4 minutes il faut décrocher QUARANTE coquilles : FOR DD 14, une "
            "action chacune. Les bernard-l'ermite ne peuvent pas le faire eux-mêmes.\n"
            "Ne coupe jamais cette séquence."
        ),
    },
    10: {
        "titre": "Les retrouvailles",
        "quand": "La toute fin. Après « Bon. » de Nine et les deux étoiles de l'inspecteur.",
        "lire": (
            "Laisse l'image, ne parle pas par-dessus.\n"
            "Puis, en dernier, la question à Bumbur : est-ce qu'il vient, ou est-ce qu'il "
            "reste ?"
        ),
        "apres": (
            "Ne réponds pas à sa place. Il faut qu'ils lui demandent, et que ça vaille le "
            "coup.\n"
            "Et les sept poulpes : personne n'a pensé à leur demander."
        ),
    },
}

FIN = {
    "titre": "FIN",
    "sous": "",
    "notes": "Écran noir (B). Laisse-les parler entre eux.",
}


# ---------------------------------------------------------------- outils

def _norm(s: str) -> str:
    return unicodedata.normalize("NFKD", s).encode("ascii", "ignore").decode().lower()


def make_slide_images() -> dict[int, Path]:
    SLIDES.mkdir(parents=True, exist_ok=True)
    found = {}
    for png in sorted(OUT.glob("*.png")):
        num = int(png.name[:2])
        dest = SLIDES / (png.stem + ".jpg")
        if not dest.exists() or dest.stat().st_mtime < png.stat().st_mtime:
            im = Image.open(png).convert("RGB")
            if im.width > TARGET_PX:
                h = round(im.height * TARGET_PX / im.width)
                im = im.resize((TARGET_PX, h), Image.LANCZOS)
            im.save(dest, quality=88, optimize=True)
        found[num] = dest
    return found


def find_audio() -> dict[int, list[Path]]:
    """Apparie les mp3 du dossier aux pistes declarees dans MUSIQUE."""
    mp3 = list(HERE.glob("*.mp3"))
    out: dict[int, list[Path]] = {}
    for num, bloc in MUSIQUE.items():
        trouves = []
        for _t, _s, _d, needle in bloc["pistes"]:
            for f in mp3:
                if needle in _norm(f.name) and f not in trouves:
                    trouves.append(f)
                    break
        if trouves:
            out[num] = trouves
    return out


def add_notes(slide, text: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(14)


def attach_audio(slide, pistes: list[Path]) -> None:
    """Une petite icone par piste, en bas a gauche, lecture au clic."""
    for i, mp3 in enumerate(pistes):
        try:
            slide.shapes.add_movie(
                str(mp3),
                Inches(0.15 + i * 0.55), Inches(SLIDE_H - 0.65),
                Inches(0.45), Inches(0.45),
                mime_type="audio/mpeg",
            )
        except Exception as exc:                  # pragma: no cover
            print(f"  ! audio non insere ({mp3.name}) : {exc}")


def text_slide(prs, titre: str, sous: str, notes: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NOIR

    box = slide.shapes.add_textbox(Inches(0), Inches(2.6), Inches(SLIDE_W), Inches(2.3))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = titre
    r = p.runs[0]
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = OR
    if sous:
        p2 = tf.add_paragraph()
        p2.text = sous
        r2 = p2.runs[0]
        r2.font.size = Pt(22)
        r2.font.color.rgb = GRIS
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
    add_notes(slide, notes)
    return slide


def image_slide(prs, img: Path, notes: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    im = Image.open(img)
    ratio_img = im.width / im.height
    ratio_slide = SLIDE_W / SLIDE_H
    if ratio_img > ratio_slide:
        h = SLIDE_H
        w = h * ratio_img
        left, top = (SLIDE_W - w) / 2, 0
    else:
        w = SLIDE_W
        h = w / ratio_img
        left, top = 0, (SLIDE_H - h) / 2
    slide.shapes.add_picture(str(img), Inches(left), Inches(top), Inches(w), Inches(h))
    add_notes(slide, notes)
    return slide


def bloc_musique(num: int, pistes: list[Path]) -> str:
    bloc = MUSIQUE[num]
    lignes = []
    for i, (titre, source, duree, needle) in enumerate(bloc["pistes"], 1):
        fichier = next((f for f in pistes if needle in _norm(f.name)), None)
        marque = " [embarqué]" if fichier else " [à télécharger]"
        lignes.append(f"   {i}. {titre} — {source}  ({duree}){marque}")
    return (
        "♪ MUSIQUE — 3 pistes, à enchaîner dans l'ordre :\n"
        + "\n".join(lignes)
        + f"\n   Départ : {bloc['cue']}\n\n"
    )


def build(with_audio: bool, dest: Path):
    images = make_slide_images()
    missing = [n for n in NOTES if n not in images]
    if missing:
        raise SystemExit(f"Images manquantes pour les scenes {missing}")
    audio = find_audio() if with_audio else {}

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    intro = (
        "les pistes présentes dans le dossier sont EMBARQUÉES. Une icône par piste en "
        "bas à gauche : clique pour lancer, reclique pour couper."
        if with_audio else
        "playlist de 30 pistes dans l'ordre, sur ton téléphone ou un second onglet."
    ) + " Tu lances toujours 2 à 4 secondes AVANT de parler, volume assez bas pour " \
        "qu'on parle normalement par-dessus."
    text_slide(prs, TITRE["titre"], TITRE["sous"],
               TITRE["notes"].replace("{MUSIQUE_INTRO}", intro))

    for num in sorted(NOTES):
        n = NOTES[num]
        pistes = audio.get(num, [])
        notes = (
            f"=== {num}. {n['titre'].upper()} ===\n"
            f"QUAND : {n['quand']}\n\n"
            f"{bloc_musique(num, pistes)}"
            f"--- À LIRE ---\n{n['lire']}\n\n"
            f"--- ENSUITE ---\n{n['apres']}\n\n"
            f"(Écran noir : touche B)"
        )
        slide = image_slide(prs, images[num], notes)
        if pistes:
            attach_audio(slide, pistes)

    text_slide(prs, FIN["titre"], FIN["sous"], FIN["notes"])
    prs.save(dest)

    mo = dest.stat().st_size / 1024 / 1024
    total = sum(len(v) for v in audio.values())
    suffixe = f", {total}/30 pistes embarquees" if audio else ""
    print(f"OK  {dest.name}  —  {len(NOTES) + 2} diapos, {mo:.1f} Mo{suffixe}")


def main() -> int:
    build(with_audio=False, dest=HERE / "Le_Service_du_soir.pptx")
    build(with_audio=True, dest=HERE / "Le_Service_du_soir_AVEC_MUSIQUE.pptx")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
